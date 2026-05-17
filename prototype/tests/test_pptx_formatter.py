"""TDD: PPT formatter renders slide JSON + MCQ JSON into a valid .pptx file."""
from pathlib import Path

from pptx import Presentation

from src.pptx_formatter import render_pptx


def test_render_pptx_creates_valid_file(tmp_path: Path):
    slides = [
        {"layout": "title", "title": "Photosynthesis", "subtitle": "Grade 7"},
        {"layout": "bullets", "title": "Process",
         "bullets": ["Sunlight", "Water", "CO2"]},
        {"layout": "two_column", "title": "In vs Out",
         "left_column": "Inputs: CO2, H2O", "right_column": "Outputs: O2, glucose"},
        {"layout": "bullets", "title": "Where it happens",
         "bullets": ["Chloroplasts in plant leaves",
                     "Mostly in palisade and spongy mesophyll cells"]},
    ]
    mcqs = [
        {"question": "What gas do plants take in?",
         "options": ["O2", "CO2", "N2", "H2"],
         "answer": "B", "explanation": "Plants absorb CO2 from the air."},
        {"question": "Where does photosynthesis happen?",
         "options": ["Roots", "Stem", "Chloroplasts", "Mitochondria"],
         "answer": "C", "explanation": "Chloroplasts contain chlorophyll."},
    ]
    out = tmp_path / "lesson.pptx"

    render_pptx(slides, mcqs, str(out))

    assert out.exists(), "pptx file was not created"
    assert out.stat().st_size > 5_000, "pptx file looks empty"

    prs = Presentation(str(out))
    assert len(prs.slides) == len(slides) + len(mcqs)
    assert prs.slides[0].shapes.title.text == "Photosynthesis"
    assert "photosynthesis happen" in prs.slides[-1].shapes.title.text.lower()


def test_render_pptx_handles_empty_mcqs(tmp_path: Path):
    slides = [{"layout": "title", "title": "Topic", "subtitle": "Sub"}]
    out = tmp_path / "no_mcqs.pptx"
    render_pptx(slides, [], str(out))
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_render_pptx_unknown_layout_falls_back_to_bullets(tmp_path: Path):
    slides = [{"layout": "weird_unknown", "title": "X", "bullets": ["a", "b"]}]
    out = tmp_path / "fallback.pptx"
    render_pptx(slides, [], str(out))
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_render_pptx_brands_title_slide_with_teacher_name(tmp_path: Path):
    slides = [{"layout": "title", "title": "Photosynthesis", "subtitle": "Grade 7"}]
    out = tmp_path / "branded.pptx"
    render_pptx(slides, [], str(out), teacher_name="Ms. Priya")
    prs = Presentation(str(out))
    subtitle_text = prs.slides[0].placeholders[1].text
    assert "Grade 7" in subtitle_text
    assert "Prepared by Ms. Priya" in subtitle_text


def test_render_pptx_brands_only_first_title_slide(tmp_path: Path):
    """If the deck has two title slides, only the first carries the byline."""
    slides = [
        {"layout": "title", "title": "Main", "subtitle": "Grade 7"},
        {"layout": "title", "title": "Part 2", "subtitle": "Continued"},
    ]
    out = tmp_path / "two_title.pptx"
    render_pptx(slides, [], str(out), teacher_name="Ms. Priya")
    prs = Presentation(str(out))
    assert "Prepared by Ms. Priya" in prs.slides[0].placeholders[1].text
    assert "Prepared by" not in prs.slides[1].placeholders[1].text


def test_render_pptx_no_teacher_name_omits_byline(tmp_path: Path):
    slides = [{"layout": "title", "title": "Photosynthesis", "subtitle": "Grade 7"}]
    out = tmp_path / "unbranded.pptx"
    render_pptx(slides, [], str(out))
    prs = Presentation(str(out))
    assert "Prepared by" not in prs.slides[0].placeholders[1].text


def test_render_pptx_applies_subject_color_to_title(tmp_path: Path):
    """Title text on the title slide should be the palette's primary color."""
    from src.theme import palette_for_subject

    slides = [{"layout": "title", "title": "Photosynthesis", "subtitle": "Grade 7"}]
    out = tmp_path / "themed.pptx"
    render_pptx(slides, [], str(out), subject="Science")

    prs = Presentation(str(out))
    expected_rgb = palette_for_subject("Science").primary
    runs = list(prs.slides[0].shapes.title.text_frame.paragraphs[0].runs)
    assert runs, "title slide title has no runs"
    color = runs[0].font.color.rgb
    assert tuple(color) == expected_rgb, (
        f"title color {tuple(color)} != palette primary {expected_rgb}"
    )


def test_render_pptx_different_subjects_get_different_palettes(tmp_path: Path):
    """Two decks with the same content but different subjects should produce
    different title colors. Sanity check that the subject argument is wired."""
    slides = [{"layout": "title", "title": "Topic", "subtitle": "Grade 7"}]

    sci_path = tmp_path / "science.pptx"
    math_path = tmp_path / "math.pptx"
    render_pptx(slides, [], str(sci_path), subject="Science")
    render_pptx(slides, [], str(math_path), subject="Mathematics")

    sci_runs = list(Presentation(str(sci_path)).slides[0].shapes.title.text_frame.paragraphs[0].runs)
    math_runs = list(Presentation(str(math_path)).slides[0].shapes.title.text_frame.paragraphs[0].runs)
    assert tuple(sci_runs[0].font.color.rgb) != tuple(math_runs[0].font.color.rgb)


def test_render_pptx_unknown_subject_uses_default_palette(tmp_path: Path):
    from src.theme import palette_for_subject
    slides = [{"layout": "title", "title": "Topic", "subtitle": "Grade 7"}]
    out = tmp_path / "default.pptx"
    render_pptx(slides, [], str(out), subject="Underwater basket weaving")
    runs = list(Presentation(str(out)).slides[0].shapes.title.text_frame.paragraphs[0].runs)
    assert tuple(runs[0].font.color.rgb) == palette_for_subject(None).primary


def _safe_autoshape_type(shape):
    """python-pptx raises ValueError if you access auto_shape_type on a
    non-autoshape (e.g., the title placeholder). Return None in that case."""
    try:
        return shape.auto_shape_type
    except (AttributeError, ValueError):
        return None


def test_render_pptx_diagram_slide_draws_boxes_and_arrows(tmp_path: Path):
    """Diagram slide: N rounded boxes + (N-1) arrows + a title."""
    from pptx.enum.shapes import MSO_SHAPE
    slides = [
        {"layout": "title", "title": "Photosynthesis", "subtitle": "Grade 7"},
        {"layout": "diagram", "title": "How it works",
         "diagram": {"kind": "process_flow", "nodes": [
             {"label": "Sunlight", "detail": "captured by chlorophyll"},
             {"label": "Leaf"},
             {"label": "Glucose"},
             {"label": "Oxygen"},
         ]}},
    ]
    out = tmp_path / "diagram.pptx"
    render_pptx(slides, [], str(out), subject="Science")
    prs = Presentation(str(out))
    assert len(prs.slides) == 2

    diagram_slide = prs.slides[1]
    shape_types = [_safe_autoshape_type(s) for s in diagram_slide.shapes]
    n_boxes = sum(1 for t in shape_types if t == MSO_SHAPE.ROUNDED_RECTANGLE)
    n_arrows = sum(1 for t in shape_types if t == MSO_SHAPE.RIGHT_ARROW)
    assert n_boxes == 4, f"expected 4 node boxes, got {n_boxes}"
    assert n_arrows == 3, f"expected 3 arrows between 4 nodes, got {n_arrows}"

    first_box = next(
        s for s in diagram_slide.shapes
        if _safe_autoshape_type(s) == MSO_SHAPE.ROUNDED_RECTANGLE
    )
    assert "Sunlight" in first_box.text_frame.text


def test_render_pptx_diagram_falls_back_to_bullets_when_nodes_empty(tmp_path: Path):
    """LLM hand-waving: layout='diagram' but no nodes. Should not crash."""
    slides = [{"layout": "diagram", "title": "How", "bullets": ["A", "B"],
               "diagram": {"kind": "process_flow", "nodes": []}}]
    out = tmp_path / "empty_diagram.pptx"
    render_pptx(slides, [], str(out))
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_render_pptx_diagram_caps_at_six_nodes(tmp_path: Path):
    """More than 6 nodes are clipped so the diagram stays legible."""
    from pptx.enum.shapes import MSO_SHAPE
    nodes = [{"label": f"Step {i}"} for i in range(1, 10)]
    slides = [{"layout": "diagram", "title": "Big diagram",
               "diagram": {"kind": "process_flow", "nodes": nodes}}]
    out = tmp_path / "cap.pptx"
    render_pptx(slides, [], str(out))
    prs = Presentation(str(out))
    n_boxes = sum(
        1 for s in prs.slides[0].shapes
        if _safe_autoshape_type(s) == MSO_SHAPE.ROUNDED_RECTANGLE
    )
    assert n_boxes == 6
