"""Render slide JSON + MCQ JSON into a .pptx file using python-pptx.

The renderer owns ALL layout primitives. Callers (the LLM upstream) only pick
a layout name and supply content — they do not control fonts, sizes, or positions.
A subject-aware color palette is applied to title text and accent shapes so the
artifact reads as intentional rather than default-themed.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from src.fonts import pptx_font_for_language
from src.theme import Palette, palette_for_subject


def render_pptx(
    slides: list[dict],
    mcqs: list[dict],
    output_path: str,
    *,
    teacher_name: str | None = None,
    subject: str | None = None,
    language: str | None = None,
) -> None:
    prs = Presentation()
    palette = palette_for_subject(subject)

    title_slide_done = False
    for spec in slides:
        layout = spec.get("layout", "bullets")
        if layout == "title":
            _add_title_slide(
                prs, spec, palette,
                teacher_name=teacher_name if not title_slide_done else None,
            )
            title_slide_done = True
        elif layout == "two_column":
            _add_two_column_slide(prs, spec, palette)
        elif layout == "diagram":
            _add_diagram_slide(prs, spec, palette)
        else:
            # bullets is the default; unknown layouts also fall through here
            _add_bullet_slide(prs, spec, palette)

    for mcq in mcqs:
        _add_mcq_slide(prs, mcq, palette)

    # Apply Tamil font (or others later) to every text run on every slide.
    # Done as a post-pass so the per-slide helpers don't each need to
    # remember to set font.name on every run they create.
    font_name = pptx_font_for_language(language)
    if font_name:
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = font_name

    prs.save(output_path)


def _rgb(palette_color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*palette_color)


def _color_title(slide, palette: Palette) -> None:
    """Recolor the title placeholder text to the palette's primary color."""
    title = slide.shapes.title
    if title is None:
        return
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = _rgb(palette.primary)


def _add_accent_bar(slide, palette: Palette) -> None:
    """Draw a thin accent rectangle along the left edge of the slide. Visual
    differentiation against generic default-themed decks."""
    from pptx.enum.shapes import MSO_SHAPE
    width = Emu(120000)   # ~0.13 in
    height = slide.part.package.presentation_part.presentation.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(palette.accent)
    shape.line.fill.background()


def _add_title_slide(
    prs: Presentation, spec: dict, palette: Palette,
    *, teacher_name: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = spec.get("title", "")
    subtitle_lines: list[str] = []
    if spec.get("subtitle"):
        subtitle_lines.append(spec["subtitle"])
    if teacher_name:
        subtitle_lines.append(f"Prepared by {teacher_name}")
    if subtitle_lines and len(slide.placeholders) > 1:
        slide.placeholders[1].text = "\n".join(subtitle_lines)
    _color_title(slide, palette)
    _add_accent_bar(slide, palette)


def _add_bullet_slide(prs: Presentation, spec: dict, palette: Palette) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = spec.get("title", "")
    body = slide.placeholders[1].text_frame
    bullets = spec.get("bullets") or []
    if not bullets:
        body.text = spec.get("body", "") or ""
    else:
        body.text = bullets[0]
        for bullet in bullets[1:]:
            p = body.add_paragraph()
            p.text = bullet
    _color_title(slide, palette)


def _add_two_column_slide(prs: Presentation, spec: dict, palette: Palette) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = spec.get("title", "")
    left = slide.placeholders[1].text_frame
    right = slide.placeholders[2].text_frame
    left.text = spec.get("left_column", "") or ""
    right.text = spec.get("right_column", "") or ""
    _color_title(slide, palette)


def _add_diagram_slide(prs: Presentation, spec: dict, palette: Palette) -> None:
    """Draw a process-flow diagram: rounded boxes connected left-to-right by
    arrows, all tinted with the subject palette. Falls back to a bullets
    slide if the diagram payload is empty (LLM hand-waving)."""
    diagram = spec.get("diagram") or {}
    nodes = diagram.get("nodes") or []
    if not nodes:
        _add_bullet_slide(prs, spec, palette)
        return

    # Layout 5 = "Title Only" — gives us the title placeholder + a blank canvas
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = spec.get("title", "")
    _color_title(slide, palette)

    n = min(len(nodes), 6)  # cap nodes for visual breathing room

    side_margin = Inches(0.5)
    diagram_top = Inches(2.5)
    node_height = Inches(1.5)
    arrow_width = Inches(0.45)
    arrow_height = Inches(0.4)

    available_width = Inches(10) - 2 * side_margin - (n - 1) * arrow_width
    node_width = available_width // n

    primary_rgb = RGBColor(*palette.primary)
    accent_rgb = RGBColor(*palette.accent)

    for i, node in enumerate(nodes[:n]):
        x = side_margin + i * (node_width + arrow_width)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, diagram_top, node_width, node_height,
        )
        box.fill.solid()
        box.fill.fore_color.rgb = accent_rgb
        box.line.color.rgb = primary_rgb
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)

        label = (node.get("label") or "").strip() or f"Step {i + 1}"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = primary_rgb

        detail = (node.get("detail") or "").strip()
        if detail:
            dp = tf.add_paragraph()
            dp.alignment = PP_ALIGN.CENTER
            dr = dp.add_run()
            dr.text = detail
            dr.font.size = Pt(10)
            dr.font.color.rgb = primary_rgb

        if i < n - 1:
            arrow_x = x + node_width
            arrow_y = diagram_top + (node_height - arrow_height) // 2
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, arrow_width, arrow_height,
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = primary_rgb
            arrow.line.fill.background()


def _add_mcq_slide(prs: Presentation, mcq: dict, palette: Palette) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = mcq["question"]
    body = slide.placeholders[1].text_frame
    options = mcq["options"]
    body.text = f"A. {options[0]}"
    for i, label in enumerate(["B", "C", "D"], start=1):
        p = body.add_paragraph()
        p.text = f"{label}. {options[i]}"
    # Answer line in the accent color so it visually separates from options
    p = body.add_paragraph()
    p.text = f"Answer: {mcq['answer']} — {mcq.get('explanation', '')}"
    for run in p.runs:
        run.font.color.rgb = _rgb(palette.primary)
        run.font.size = Pt(14)
    _color_title(slide, palette)
